"""
Export Service for ValtricAI
Generates Excel, PDF, and PowerPoint files from structured data
"""

import io
import os
from typing import Dict, Any, List, Optional
from datetime import datetime
import uuid
import tempfile
from pathlib import Path

# Excel generation
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import LineChart, Reference

# PDF generation
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class ExportService:
    """Service for exporting data to various formats"""
    
    def __init__(self):
        # Create temp directory for exports
        self.export_dir = Path(tempfile.gettempdir()) / "valtric_exports"
        self.export_dir.mkdir(exist_ok=True)
        
        # Store generated files temporarily (in production, use Redis/S3)
        self.file_cache = {}
    
    def create_excel_file(self, data: Dict[str, Any], session_id: str) -> str:
        """Create Excel file from structured data"""
        wb = openpyxl.Workbook()
        ws = wb.active
        
        # Set title
        ws.title = data.get('type', 'Data Export')
        
        # Style configuration
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center")
        
        if 'headers' in data and 'rows' in data:
            # Write headers
            headers = data['headers']
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
            
            # Write data rows
            for row_idx, row_data in enumerate(data['rows'], 2):
                for col_idx, header in enumerate(headers, 1):
                    value = row_data.get(header, '')
                    # Convert to number if possible
                    try:
                        if isinstance(value, str) and value.replace('.', '').replace('-', '').isdigit():
                            value = float(value)
                    except:
                        pass
                    ws.cell(row=row_idx, column=col_idx, value=value)
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # Add chart if financial data
            if data.get('type') == 'financial' and len(data['rows']) > 1:
                chart = LineChart()
                chart.title = "Financial Projection"
                chart.y_axis.title = "Amount ($)"
                chart.x_axis.title = "Period"
                
                # Find revenue and profit columns
                revenue_col = headers.index('Revenue') + 1 if 'Revenue' in headers else None
                profit_col = headers.index('Profit') + 1 if 'Profit' in headers else None
                
                if revenue_col:
                    data_ref = Reference(ws, min_col=revenue_col, min_row=1, 
                                        max_row=len(data['rows']) + 1)
                    chart.add_data(data_ref, titles_from_data=True)
                    
                    categories = Reference(ws, min_col=1, min_row=2, 
                                         max_row=len(data['rows']) + 1)
                    chart.set_categories(categories)
                    
                    ws.add_chart(chart, f"A{len(data['rows']) + 4}")
        
        elif data.get('type') == 'swot' and 'data' in data:
            # SWOT specific layout
            swot_data = data['data']
            
            # Title
            ws.merge_cells('A1:D1')
            title_cell = ws['A1']
            title_cell.value = "SWOT Analysis"
            title_cell.font = Font(bold=True, size=16)
            title_cell.alignment = Alignment(horizontal="center")
            
            # SWOT quadrants
            quadrants = [
                ('A3', 'Strengths', swot_data.get('strengths', [])),
                ('C3', 'Weaknesses', swot_data.get('weaknesses', [])),
                ('A15', 'Opportunities', swot_data.get('opportunities', [])),
                ('C15', 'Threats', swot_data.get('threats', []))
            ]
            
            for start_cell, title, items in quadrants:
                # Quadrant title
                ws[start_cell] = title
                ws[start_cell].font = Font(bold=True, size=12)
                ws[start_cell].fill = header_fill
                ws[start_cell].font = header_font
                
                # Items
                row_offset = 1
                for item in items:
                    item_text = item['item'] if isinstance(item, dict) else item
                    cell_ref = ws.cell(row=ws[start_cell].row + row_offset, 
                                      column=ws[start_cell].column, 
                                      value=f"• {item_text}")
                    row_offset += 1
        
        # Save file
        file_id = f"{session_id}_{uuid.uuid4().hex[:8]}"
        file_path = self.export_dir / f"{file_id}.xlsx"
        wb.save(file_path)
        
        # Cache file path
        self.file_cache[file_id] = {
            'path': str(file_path),
            'type': 'excel',
            'created_at': datetime.now()
        }
        
        return file_id
    
    def create_pdf_report(self, data: Dict[str, Any], session_id: str) -> str:
        """Create PDF report from structured data"""
        file_id = f"{session_id}_{uuid.uuid4().hex[:8]}"
        file_path = self.export_dir / f"{file_id}.pdf"
        
        # Create PDF document
        doc = SimpleDocTemplate(str(file_path), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Title style
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1e40af'),
            spaceAfter=30,
            alignment=1  # Center
        )
        
        # Add title
        title_text = f"ValtricAI {data.get('type', 'Report').title()} Analysis"
        story.append(Paragraph(title_text, title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # Add metadata
        metadata_text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        story.append(Paragraph(metadata_text, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Add data content
        if 'headers' in data and 'rows' in data:
            # Create table
            table_data = [data['headers']]
            for row in data['rows']:
                table_row = [str(row.get(h, '')) for h in data['headers']]
                table_data.append(table_row)
            
            # Create table with styling
            table = Table(table_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1e40af')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
            ]))
            story.append(table)
            
        elif data.get('type') == 'swot' and 'data' in data:
            # SWOT specific layout
            swot_data = data['data']
            
            for category in ['strengths', 'weaknesses', 'opportunities', 'threats']:
                story.append(Paragraph(category.title(), styles['Heading2']))
                items = swot_data.get(category, [])
                for item in items:
                    item_text = item['item'] if isinstance(item, dict) else item
                    story.append(Paragraph(f"• {item_text}", styles['Normal']))
                story.append(Spacer(1, 0.2*inch))
        
        # Add footer
        story.append(PageBreak())
        footer_text = "Generated by ValtricAI Consulting Agent"
        story.append(Paragraph(footer_text, styles['Normal']))
        
        # Build PDF
        doc.build(story)
        
        # Cache file
        self.file_cache[file_id] = {
            'path': str(file_path),
            'type': 'pdf',
            'created_at': datetime.now()
        }
        
        return file_id
    
    def create_powerpoint(self, data: Dict[str, Any], session_id: str) -> str:
        """Create PowerPoint presentation from structured data"""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"{data.get('type', 'Analysis').title()} Report"
        subtitle.text = f"Generated by ValtricAI | {datetime.now().strftime('%B %d, %Y')}"
        
        # Content slides based on data type
        if data.get('type') == 'swot' and 'data' in data:
            # SWOT Analysis slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            title_shape.text = "SWOT Analysis"
            
            # Create 2x2 table for SWOT
            rows, cols = 2, 2
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(4)
            
            table = shapes.add_table(rows, cols, left, top, width, height).table
            
            # SWOT categories
            swot_data = data['data']
            categories = [
                ('Strengths', swot_data.get('strengths', [])),
                ('Weaknesses', swot_data.get('weaknesses', [])),
                ('Opportunities', swot_data.get('opportunities', [])),
                ('Threats', swot_data.get('threats', []))
            ]
            
            for idx, (category, items) in enumerate(categories):
                row = idx // 2
                col = idx % 2
                cell = table.cell(row, col)
                
                # Add category title
                text_frame = cell.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = category
                p.font.bold = True
                p.font.size = Pt(14)
                
                # Add items
                for item in items[:3]:  # Limit to 3 items per quadrant
                    p = text_frame.add_paragraph()
                    item_text = item['item'] if isinstance(item, dict) else item
                    p.text = f"• {item_text}"
                    p.font.size = Pt(11)
                    p.level = 1
        
        elif 'headers' in data and 'rows' in data:
            # Data table slide
            bullet_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title = shapes.title
            title.text = f"{data.get('type', 'Data').title()} Analysis"
            
            # Add table
            rows = min(len(data['rows']) + 1, 10)  # Limit rows for readability
            cols = len(data['headers'])
            
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(4)
            
            table = shapes.add_table(rows, cols, left, top, width, height).table
            
            # Add headers
            for col_idx, header in enumerate(data['headers']):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
            
            # Add data
            for row_idx, row_data in enumerate(data['rows'][:rows-1], 1):
                for col_idx, header in enumerate(data['headers']):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(row_data.get(header, ''))
        
        # Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Key Takeaways"
        
        tf = body_shape.text_frame
        tf.text = "• Data-driven insights generated by AI analysis"
        
        p = tf.add_paragraph()
        p.text = "• Structured framework for decision making"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "• Actionable recommendations based on analysis"
        p.level = 0
        
        # Save file
        file_id = f"{session_id}_{uuid.uuid4().hex[:8]}"
        file_path = self.export_dir / f"{file_id}.pptx"
        prs.save(file_path)
        
        # Cache file
        self.file_cache[file_id] = {
            'path': str(file_path),
            'type': 'powerpoint',
            'created_at': datetime.now()
        }
        
        return file_id
    
    def get_file_path(self, file_id: str) -> Optional[str]:
        """Get file path from cache"""
        file_info = self.file_cache.get(file_id)
        if file_info and os.path.exists(file_info['path']):
            return file_info['path']
        return None
    
    def cleanup_old_files(self, max_age_hours: int = 24):
        """Clean up old export files"""
        current_time = datetime.now()
        for file_id, file_info in list(self.file_cache.items()):
            age = current_time - file_info['created_at']
            if age.total_seconds() > max_age_hours * 3600:
                # Delete file
                try:
                    os.remove(file_info['path'])
                except:
                    pass
                # Remove from cache
                del self.file_cache[file_id]

# Global instance
export_service = ExportService()